#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
대우오피스 AIoT 기반 무인 공유오피스 고도화 사업계획서
Word(.docx) + PowerPoint(.pptx) 생성 스크립트
"""

import os
from docx import Document
from docx.shared import Pt, RGBColor, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as PPTRGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ─────────────────────────────────────────────
# WORD 문서 생성
# ─────────────────────────────────────────────
def create_word():
    doc = Document()

    # 기본 폰트 설정
    style = doc.styles['Normal']
    style.font.name = '맑은 고딕'
    style.font.size = Pt(10.5)

    def add_heading(text, level=1, color=None):
        p = doc.add_heading(text, level=level)
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
        run = p.runs[0] if p.runs else p.add_run(text)
        run.font.name = '맑은 고딕'
        if color:
            run.font.color.rgb = RGBColor(*color)
        return p

    def add_para(text, bold=False, indent=False, bullet=False):
        p = doc.add_paragraph()
        if indent:
            p.paragraph_format.left_indent = Cm(0.8)
        if bullet:
            p.style = doc.styles['List Bullet']
        run = p.add_run(text)
        run.font.name = '맑은 고딕'
        run.font.size = Pt(10.5)
        run.bold = bold
        return p

    def add_table_row(table, cells):
        row = table.add_row()
        for i, text in enumerate(cells):
            cell = row.cells[i]
            cell.text = text
            for run in cell.paragraphs[0].runs:
                run.font.name = '맑은 고딕'
                run.font.size = Pt(10)

    # ── 표지 ──
    doc.add_paragraph()
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_para.add_run('사  업  계  획  서')
    title_run.font.name = '맑은 고딕'
    title_run.font.size = Pt(28)
    title_run.bold = True
    title_run.font.color.rgb = RGBColor(0x1e, 0x40, 0xaf)  # blue-800

    doc.add_paragraph()
    sub_para = doc.add_paragraph()
    sub_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub_run = sub_para.add_run('AIoT 기반 무인 공유오피스 고도화 사업')
    sub_run.font.name = '맑은 고딕'
    sub_run.font.size = Pt(16)
    sub_run.bold = True

    doc.add_paragraph()
    doc.add_paragraph()

    info_para = doc.add_paragraph()
    info_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    info_run = info_para.add_run(
        '사업자명: 대우오피스\n'
        '대표자: (대표자명)\n'
        '사업자등록번호: (사업자번호)\n'
        '소재지: (사업장 주소)\n'
        '작성일: 2026년 2월'
    )
    info_run.font.name = '맑은 고딕'
    info_run.font.size = Pt(12)

    doc.add_page_break()

    # ── 목차 ──
    add_heading('목   차', level=1, color=(0x1e, 0x40, 0xaf))
    toc_items = [
        ('1장', '기업 개요 및 사업 비전', '3'),
        ('2장', '시장 분석 및 문제 제기', '5'),
        ('3장', '핵심 기술: AIoT 기반 무인 운영 시스템', '7'),
        ('4장', '공간 운영 전략: 하이브리드 스페이스', '10'),
        ('5장', '소요 자금 및 조달 계획', '12'),
        ('6장', '기대 효과 및 성장 전략', '14'),
        ('별첨', '공간 사용 승인 관련 서류 목록', '16'),
    ]
    toc_table = doc.add_table(rows=1, cols=3)
    toc_table.style = 'Table Grid'
    hdr = toc_table.rows[0].cells
    hdr[0].text = '구분'
    hdr[1].text = '항목'
    hdr[2].text = '페이지'
    for item in toc_items:
        add_table_row(toc_table, item)

    doc.add_page_break()

    # ── 1장 ──
    add_heading('제 1장. 기업 개요 및 사업 비전', level=1, color=(0x1e, 0x40, 0xaf))

    add_heading('1-1. 기업 소개', level=2)
    add_para('대우오피스는 2014년 11월 설립되어 약 11년의 업력을 보유한 지역 거점 비즈니스 센터입니다.')
    add_para('현재 경기도 일산에 위치한 집합상가 건물 내에서 1인 기업 형태로 소호사무실을 운영 중이며, '
             '입주 기업들에게 안정적이고 합리적인 업무 공간을 제공해왔습니다.', indent=True)
    add_para('영업 실평수 75평의 스 군룰 소호사무실(20실)과 함께, 최근 건물 공용부 내 16평의 커뮤니티 라운지 공간을 '
             '추가 확보하여 코워킹 오피스로의 서비스 고도화를 추진 중입니다.', indent=True)

    add_heading('1-2. 대표자 전문성', level=2)
    bullets_1_2 = [
        '부동산 임대업 11년 차 운영 전문가 (1인 기업 운영 실적 보유)',
        '집합상가 관리단 총무 3연임(12·13·14대) — 건물 운영 및 법률적 이해도 상위',
        'Google AI 및 생성형 AI 기술 활용 역량 적극 습득 중 (Coursera 수강)',
        '임대 관리 자동화 웹 서비스(Nabido) 직접 개발 진행 중',
    ]
    for b in bullets_1_2:
        add_para(b, bullet=True)

    add_heading('1-3. 사업 전환(Pivot) 배경 및 비전', level=2)
    add_para('기존 소호사무실 모델은 공간 임대에 그치는 수동적 수익 구조였습니다. 이에 AI 기술과 IoT 인프라를 결합하여 '
             '"1인 무인 운영이 가능한 스마트 공유오피스"로 사업을 고도화합니다.')
    add_para('[비전] 대우오피스를 소규모 임대 사업자의 AI 전환 성공 모델로 구축하고, '
             '이 시스템을 외부에 솔루션화하는 B2B 플랫폼으로 확장한다.', bold=True, indent=True)

    doc.add_page_break()

    # ── 2장 ──
    add_heading('제 2장. 시장 분석 및 문제 제기', level=1, color=(0x1e, 0x40, 0xaf))

    add_heading('2-1. 시장 환경', level=2)
    add_para('2024~2026년 국내 1인 창조기업 및 프리랜서 증가로 유연한 소규모 업무 공간 수요가 지속 확대되고 있습니다.')
    market_table = doc.add_table(rows=4, cols=2)
    market_table.style = 'Table Grid'
    for row, data in zip(market_table.rows, [
        ('지표', '현황'),
        ('국내 공유오피스 시장 규모', '2025년 약 1.2조 원 (年 15% 성장)'),
        ('1인 창조기업 등록 수', '약 22만 개 사 (2025년 기준)'),
        ('스마트오피스 도입 의향 비율', '중소기업 64% 이상 관심 표명'),
    ]):
        for i, text in enumerate(data):
            row.cells[i].text = text

    add_heading('2-2. 기존 문제점', level=2)
    problems = [
        '에너지 낭비: 중앙 냉난방 방식으로 공실 룸에도 전기 소모 지속',
        '인건비 부담: 관리자 상주 없이는 보안·출입 통제 어려움',
        '커뮤니티 부재: 입주사 간 네트워킹 공간 없어 이탈률 증가',
        '수작업 정산: 월세·전기료 개별 청구에 수십 시간 소요',
    ]
    for p in problems:
        add_para(p, bullet=True)

    add_heading('2-3. 해결 솔루션', level=2)
    add_para('"초저비용 고효율 AIoT 공유오피스" — AI가 관리하고, 사람은 전략에 집중한다.', bold=True)

    doc.add_page_break()

    # ── 3장 ──
    add_heading('제 3장. 핵심 기술: AIoT 기반 무인 운영 시스템', level=1, color=(0x1e, 0x40, 0xaf))

    add_heading('3-1. 스마트 에너지 관리 시스템 (EMS)', level=2)
    ems_items = [
        '재실 감지 센서(PIR/밀리미터파): 20개 룸 각각에 설치, 사람 없을 시 자동 절전 전환',
        '대형 멀티형 실외기 1대 + 개별 실내기 20대 시스템 — 중앙 제어로 에너지 최적화',
        'AI 알고리즘으로 온도 패턴 학습 → 계절별 자동 쾌적 온도 유지',
        '전기 사용량 실시간 측정 → 월말 자동 청구서 생성 (수작업 000% 제거)',
    ]
    for item in ems_items:
        add_para(item, bullet=True)

    add_heading('3-2. 무인 관제 시스템', level=2)
    control_items = [
        '모바일 앱/웹 대시보드: 전체 20개 룸의 온도·습도·전력·보안 상태 원격 확인',
        '스마트 출입 통제: QR코드 또는 NFC 기반 무인 출입 인증',
        'AI 챗봇: 입주사 민원·계약 문의·관리비 조회 자동 응대 (24/7)',
        '이상 감지 알림: 화재·침수·침입 센서 연동, 실시간 SMS·카카오톡 알림',
    ]
    for item in control_items:
        add_para(item, bullet=True)

    add_heading('3-3. 1인 운영 최적화 (운영 자동화)', level=2)
    auto_items = [
        '임대 관리 AI 솔루션(Nabido): 계약 만료 D-30 자동 안내, 미납 자동 독촉',
        '세금계산서 자동 발행: 팝빌 API 연동, 수납 완료 즉시 자동 발행',
        '카카오 알림톡 자동화: 납부일 D-3 사전 안내 → D-Day 청구 → D+1 독촉',
        '1인 대표자가 스마트폰 하나로 전체 오피스 운영 — 상주 직원 不필요',
    ]
    for item in auto_items:
        add_para(item, bullet=True)

    doc.add_page_break()

    # ── 4장 ──
    add_heading('제 4장. 공간 운영 전략: 하이브리드 스페이스', level=1, color=(0x1e, 0x40, 0xaf))

    add_heading('4-1. Private Zone — 자가 소유 20개 독립 룸 (75평)', level=2)
    add_para('현황: 완전 자가 소유 공간으로 안정적 임대 수익 확보 중 (공실률 최소화)')
    add_para('고도화 계획:')
    private_items = [
        '노후 개별 에어컨 → 신형 시스템 에어컨(AI 제어 가능 모델) 전면 교체',
        '각 룸에 재실 감지 센서 + 스마트 콘센트 설치',
        '도어락 → NFC 스마트 잠금 장치 교체 (입주사 앱 제어)',
    ]
    for item in private_items:
        add_para(item, bullet=True)

    add_heading('4-2. Public Zone — 커뮤니티 라운지 (16평)', level=2)
    add_para('공간 확보 현황:')
    public_items = [
        '입지: 집합상가(3,800평) 건물 공용부 16평 — 월 이용료 20만 원으로 전용 사용',
        '법적 근거: 관리단 총무(3연임) 자격으로 관리단 회의를 통해 정식 전용 사용 승인 획득',
        '안정성: 관리단 결의 및 사용 계약서(장기) 체결로 최소 3년 이상 사용 보장',
    ]
    for item in public_items:
        add_para(item, bullet=True)

    add_para('활용 계획:')
    lounge_items = [
        '24시간 무인 커뮤니티 라운지 — 입주사 교류·미팅·협업 공간',
        '이동식 방음 부스 + 모듈형 회의 테이블: 고정 공사 없이 비품으로 공간 조성',
        '빔프로젝터·디지털 화이트보드 설치로 소규모 세미나·워크숍 가능',
        '무인 키오스크로 라운지 사용 예약·출입 셀프 관리',
    ]
    for item in lounge_items:
        add_para(item, bullet=True)

    doc.add_page_break()

    # ── 5장 ──
    add_heading('제 5장. 소요 자금 및 조달 계획', level=1, color=(0x1e, 0x40, 0xaf))

    add_heading('5-1. 총 소요 자금: 약 75,000,000원', level=2)

    budget_table = doc.add_table(rows=1, cols=4)
    budget_table.style = 'Table Grid'
    headers = budget_table.rows[0].cells
    for i, h in enumerate(['항목', '상세 내용', '추정 금액', '비고']):
        headers[i].text = h
        for run in headers[i].paragraphs[0].runs:
            run.bold = True

    budget_data = [
        ('스마트 냉난방 시스템', '실외기(멀티V 대형) 1대\n+ 실내기 20대 + 배관 공사', '35,000,000 ~ 45,000,000원', '시설 투자 / 핵심 자산'),
        ('AIoT 관제 시스템', '재실감지 센서 20개\n스마트 도어락 20개\n전력 측정기·컨트롤러\n앱/웹 개발비', '10,000,000 ~ 15,000,000원', '소프트웨어 + 하드웨어'),
        ('라운지 비품', '이동식 방음 부스\n회의 테이블·의자\n빔프로젝터·화이트보드', '5,000,000 ~ 10,000,000원', '이동 가능 자산'),
        ('무인 운영 솔루션', '키오스크 1대\n출입통제 시스템\nAI 앱 고도화', '5,000,000원', '운영 자동화'),
        ('합 계', '', '약 55,000,000 ~ 75,000,000원', '정책자금 7,000만 + 자부담'),
    ]
    for row_data in budget_data:
        add_table_row(budget_table, row_data)

    add_heading('5-2. 자금 조달 방안', level=2)
    fund_items = [
        '[소상공인 정책자금] 일반경영안정자금 (대리대출): 70,000,000원 — 시설 교체 + 시스템 구축 운전자금 명목',
        '[자부담] 자기 자금: 5,000,000원',
        '[보조금] 스마트상점 기술보급사업 (3월 매년 공고): 키오스크·스마트기기 최대 700만 원 별도 지원 신청 예정',
    ]
    for item in fund_items:
        add_para(item, bullet=True)

    add_heading('5-3. 정책자금 신청 전략', level=2)
    strategy_items = [
        '신청 명목: "인테리어"가 아닌 "스마트 시스템 도입에 따른 운전자금" — 시설 자금보다 승인율 높음',
        '업종 예외 적용: 701206(공유오피스) — 부동산업 예외 지원 대상 명시적 소명',
        '공간 안정성 증빙: 관리단 총무 3연임 이력 및 전용 사용 계약서(별첨) 제출',
        '기술 혁신성 어필: AI 시스템을 통한 에너지 절감·무인화 계획 구체적 제시',
    ]
    for item in strategy_items:
        add_para(item, bullet=True)

    doc.add_page_break()

    # ── 6장 ──
    add_heading('제 6장. 기대 효과 및 성장 전략', level=1, color=(0x1e, 0x40, 0xaf))

    add_heading('6-1. 재무적 효과', level=2)
    effect_table = doc.add_table(rows=1, cols=3)
    effect_table.style = 'Table Grid'
    for i, h in enumerate(['항목', '현재', '도입 후 목표']):
        effect_table.rows[0].cells[i].text = h
    for row_data in [
        ('상주 관리 인건비', '월 250만 원', '월 0원 (무인화)'),
        ('전기료 (냉난방)', '월 150만 원 (추정)', '월 100만 원 (33% 절감)'),
        ('미납 독촉 업무', '월 20시간 이상', '완전 자동화'),
        ('세금계산서 발행', '수작업 월 3시간', 'API 자동 발행'),
        ('연간 순수익 증가', '—', '+3,000만 원 이상'),
    ]:
        add_table_row(effect_table, row_data)

    add_heading('6-2. 사업 확장 로드맵', level=2)
    roadmap_items = [
        '[2026년] 대우오피스 내 AIoT 시스템 1차 구축 완료 — 자체 검증',
        '[2027년] 개발 시스템을 SaaS 형태로 외부 소규모 임대 사업자에게 공개 (Nabido 플랫폼)',
        '[2028년] 타 지역 소호사무실·공유오피스 50개 이상 적용 — B2B 솔루션 수익화',
    ]
    for item in roadmap_items:
        add_para(item, bullet=True)

    add_heading('6-3. 사회적 가치', level=2)
    social_items = [
        '1인 창조기업·스타트업의 저렴하고 스마트한 창업 공간 제공',
        '소규모 임대 사업의 디지털 전환 성공 모델 확립',
        '에너지 절감을 통한 탄소 중립 기여',
    ]
    for item in social_items:
        add_para(item, bullet=True)

    doc.add_page_break()

    # ── 별첨 ──
    add_heading('별첨. 제출 서류 목록', level=1, color=(0x1e, 0x40, 0xaf))
    annex_items = [
        '사업자등록증 사본',
        '집합상가 관리단 전용 사용 승인 결의서 (관리단 회의록 사본)',
        '공용부분 전용 사용 계약서 (관리단 총무 직인)',
        '건물 관리단 총무 재직 증명서류',
        '대우오피스 공간 사진 (전체 내부, 라운지 공간)',
        'AIoT 시스템 구축 견적서 (냉난방, 센서, 소프트웨어)',
        '소상공인 확인서',
        '최근 3개월 부가세 신고서 또는 매출 증빙',
    ]
    for item in annex_items:
        add_para(f'□  {item}')

    # 저장
    output_path = os.path.join(OUTPUT_DIR, '대우오피스_사업계획서.docx')
    doc.save(output_path)
    print(f'✅ Word 저장 완료: {output_path}')
    return output_path


# ─────────────────────────────────────────────
# PPT 생성
# ─────────────────────────────────────────────
def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    BLUE = PPTRGBColor(0x1e, 0x40, 0xaf)
    BLUE_LIGHT = PPTRGBColor(0xdb, 0xe9, 0xfe)
    WHITE = PPTRGBColor(0xff, 0xff, 0xff)
    DARK = PPTRGBColor(0x0f, 0x17, 0x2a)
    GRAY = PPTRGBColor(0x64, 0x74, 0x8b)
    EMERALD = PPTRGBColor(0x05, 0x96, 0x69)

    blank_layout = prs.slide_layouts[6]  # 완전 빈 레이아웃

    def add_slide():
        return prs.slides.add_slide(blank_layout)

    def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
        shape = slide.shapes.add_shape(
            1, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        if line_color:
            shape.line.color.rgb = line_color
        else:
            shape.line.fill.background()
        return shape

    def add_text_box(slide, text, left, top, width, height,
                     font_size=20, bold=False, color=None, align=PP_ALIGN.LEFT,
                     font_name='맑은 고딕', wrap=True):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.name = font_name
        if color:
            run.font.color.rgb = color
        return txBox

    def slide_header(slide, title, subtitle=None):
        """왼쪽 파란 세로 바 + 제목"""
        add_rect(slide, 0, 0, 0.25, 7.5, fill_color=BLUE)
        add_text_box(slide, title, 0.5, 0.2, 12, 0.7,
                     font_size=28, bold=True, color=BLUE)
        if subtitle:
            add_text_box(slide, subtitle, 0.5, 0.85, 12, 0.4,
                         font_size=13, color=GRAY)
        # 구분선
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(12.5), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = BLUE_LIGHT
        line.line.fill.background()

    # ── 슬라이드 1: 표지 ──
    slide1 = add_slide()
    add_rect(slide1, 0, 0, 13.33, 7.5, fill_color=DARK)
    add_rect(slide1, 0, 5.5, 13.33, 2, fill_color=BLUE)
    add_text_box(slide1, '사 업 계 획 서', 1, 1.2, 11, 1.5,
                 font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide1, 'AIoT 기반 무인 공유오피스 고도화 사업', 1, 2.8, 11, 0.9,
                 font_size=22, bold=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
    add_text_box(slide1, '대우오피스  |  2026년 2월', 1, 5.7, 11, 0.6,
                 font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

    # ── 슬라이드 2: 목차 ──
    slide2 = add_slide()
    add_rect(slide2, 0, 0, 13.33, 7.5, fill_color=PPTRGBColor(0xf8, 0xfa, 0xfc))
    slide_header(slide2, '목  차')
    chapters = [
        ('01', '기업 개요 및 사업 비전'),
        ('02', '시장 분석 및 문제 제기'),
        ('03', '핵심 기술: AIoT 무인 운영 시스템'),
        ('04', '공간 운영 전략: 하이브리드 스페이스'),
        ('05', '소요 자금 및 조달 계획'),
        ('06', '기대 효과 및 성장 전략'),
    ]
    for i, (num, title) in enumerate(chapters):
        col = i % 2
        row = i // 2
        x = 0.7 + col * 6.3
        y = 1.4 + row * 1.6
        add_rect(slide2, x, y, 5.8, 1.3, fill_color=WHITE, line_color=BLUE_LIGHT)
        add_text_box(slide2, num, x + 0.15, y + 0.1, 0.8, 0.5,
                     font_size=22, bold=True, color=BLUE)
        add_text_box(slide2, title, x + 0.9, y + 0.15, 4.7, 0.9,
                     font_size=14, bold=True, color=DARK)

    # ── 슬라이드 3: 기업 개요 ──
    slide3 = add_slide()
    slide_header(slide3, '01. 기업 개요 및 사업 비전', '14년 업력의 지역 거점 비즈니스 센터 — AI로 새롭게 도약')

    cards = [
        ('🏢 기업 현황', '2014년 11월 설립\n운영 11년 차 (1인 기업)\n경기도 일산 위치\n75평 소호사무실 + 16평 라운지'),
        ('👤 대표자 전문성', '부동산 임대업 11년 전문가\n관리단 총무 3연임 (법적 이해 상위)\nGoogle AI 기술 습득 중\n임대관리 웹서비스 직접 개발 중'),
        ('🚀 사업 비전', '단순 임대업 →\nAI 무인 관리 플랫폼으로 Pivot\n1인 운영 완전 자동화 실현\n솔루션 외부 판매(SaaS) 목표'),
    ]
    for i, (title, body) in enumerate(cards):
        x = 0.5 + i * 4.2
        add_rect(slide3, x, 1.4, 3.9, 5.5, fill_color=WHITE, line_color=BLUE_LIGHT)
        add_rect(slide3, x, 1.4, 3.9, 0.65, fill_color=BLUE)
        add_text_box(slide3, title, x + 0.15, 1.45, 3.6, 0.55,
                     font_size=14, bold=True, color=WHITE)
        add_text_box(slide3, body, x + 0.15, 2.2, 3.6, 4.5,
                     font_size=12, color=DARK)

    # ── 슬라이드 4: 시장 분석 ──
    slide4 = add_slide()
    slide_header(slide4, '02. 시장 분석 및 문제 제기', '성장하는 시장, 그러나 기존 방식으로는 한계')

    add_rect(slide4, 0.5, 1.4, 5.8, 2.5, fill_color=PPTRGBColor(0xef, 0xf6, 0xff), line_color=BLUE_LIGHT)
    add_text_box(slide4, '📊 시장 환경', 0.65, 1.45, 5, 0.5, font_size=14, bold=True, color=BLUE)
    add_text_box(slide4, '• 공유오피스 시장 규모: 1.2조 원 (年 15% 성장)\n• 1인 창조기업 등록: 약 22만 개 사\n• 스마트오피스 도입 의향: 중소기업 64%', 0.65, 1.95, 5.5, 1.8, font_size=12, color=DARK)

    add_rect(slide4, 0.5, 4.1, 5.8, 2.8, fill_color=PPTRGBColor(0xff, 0xf1, 0xf2), line_color=PPTRGBColor(0xfb, 0xca, 0xca))
    add_text_box(slide4, '⚠️ 기존 소호사무실의 문제점', 0.65, 4.15, 5, 0.5, font_size=14, bold=True, color=PPTRGBColor(0xbe, 0x12, 0x3c))
    add_text_box(slide4, '• 중앙 냉난방 → 에너지 낭비 심각\n• 관리자 상주 불가 → 보안 취약\n• 커뮤니티 부재 → 입주사 이탈률 증가\n• 수작업 정산 → 월 20시간 이상 낭비', 0.65, 4.65, 5.5, 2.1, font_size=12, color=DARK)

    add_rect(slide4, 6.8, 1.4, 6.0, 5.5, fill_color=BLUE)
    add_text_box(slide4, '💡 우리의 해결책', 7.0, 1.5, 5.5, 0.55, font_size=16, bold=True, color=WHITE)
    add_text_box(slide4, 'AI + IoT 결합\n\n"초저비용 고효율\n무인 공유오피스"\n\n관리 인건비 ZERO\n에너지 30% 절감\n자동화 정산 시스템',
                 7.0, 2.1, 5.5, 4.5, font_size=14, color=WHITE)

    # ── 슬라이드 5: AIoT 기술 ──
    slide5 = add_slide()
    slide_header(slide5, '03. 핵심 기술: AIoT 무인 운영 시스템', '1인 대표자가 스마트폰 하나로 전체 오피스 운영')

    tech_items = [
        ('⚡ 스마트 에너지 관리', '• 재실 감지 센서 → 자동 절전\n• 멀티형 시스템 에어컨 중앙 제어\n• AI 온도 패턴 학습'),
        ('🔐 무인 출입 관제', '• NFC/QR 스마트 도어락\n• 모바일 앱 원격 모니터링\n• 이상 감지 실시간 알림'),
        ('💰 자동 정산 시스템', '• 전기 사용량 자동 측정\n• 월별 관리비 AI 자동 청구\n• 세금계산서 자동 발행'),
        ('🤖 AI 챗봇', '• 24/7 입주사 민원 자동 응대\n• 계약 만료 자동 안내\n• 미납 독촉 카카오 알림톡'),
    ]
    for i, (title, body) in enumerate(tech_items):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.3
        y = 1.4 + row * 2.8
        add_rect(slide5, x, y, 6.0, 2.5, fill_color=WHITE, line_color=BLUE_LIGHT)
        add_text_box(slide5, title, x + 0.2, y + 0.15, 5.5, 0.55, font_size=14, bold=True, color=BLUE)
        add_text_box(slide5, body, x + 0.2, y + 0.7, 5.5, 1.7, font_size=12, color=DARK)

    # ── 슬라이드 6: 공간 전략 ──
    slide6 = add_slide()
    slide_header(slide6, '04. 공간 운영 전략: 하이브리드 스페이스', '자가 소유(Private) + 라운지(Public) 최적 결합')

    add_rect(slide6, 0.5, 1.4, 5.8, 5.5, fill_color=PPTRGBColor(0xef, 0xf6, 0xff), line_color=BLUE_LIGHT)
    add_rect(slide6, 0.5, 1.4, 5.8, 0.65, fill_color=BLUE)
    add_text_box(slide6, '🔒 PRIVATE ZONE (20룸 / 75평)', 0.65, 1.45, 5.5, 0.55, font_size=13, bold=True, color=WHITE)
    add_text_box(slide6, '• 완전 자가 소유 — 안정적 임대 수익\n• 시스템 에어컨 AI 제어 교체\n• 룸별 재실 감지 + 전력 측정 설치\n• NFC 스마트 도어락 전면 교체',
                 0.65, 2.15, 5.5, 4.5, font_size=12, color=DARK)

    add_rect(slide6, 6.7, 1.4, 6.1, 5.5, fill_color=PPTRGBColor(0xf0, 0xfd, 0xf4), line_color=PPTRGBColor(0xa7, 0xf3, 0xd0))
    add_rect(slide6, 6.7, 1.4, 6.1, 0.65, fill_color=EMERALD)
    add_text_box(slide6, '🌐 PUBLIC ZONE (라운지 / 16평)', 6.85, 1.45, 5.8, 0.55, font_size=13, bold=True, color=WHITE)
    add_text_box(slide6, '• 집합상가 공용부 전용 사용 (월 20만 원)\n• 관리단 총무 3연임 → 장기 사용권 안정 확보\n• 이동식 방음 부스 + 모듈형 회의 테이블\n• 24시간 무인 운영 (키오스크 + 스마트 예약)\n• 입주사 간 네트워킹·협업 거점',
                 6.85, 2.15, 5.8, 4.5, font_size=12, color=DARK)

    # ── 슬라이드 7: 자금 계획 ──
    slide7 = add_slide()
    slide_header(slide7, '05. 소요 자금 및 조달 계획', '총 약 75,000,000원 / 소상공인 정책자금 70,000,000원 신청')

    budget_data_ppt = [
        ('스마트 냉난방 시스템', '35~45,000,000원', '시설 투자 핵심'),
        ('AIoT 관제 시스템 (소프트웨어 + 센서)', '10~15,000,000원', '기술 혁신 핵심'),
        ('라운지 비품 (이동식 부스, 가구 등)', '5~10,000,000원', '이동 가능 자산'),
        ('무인 운영 솔루션 (키오스크, AI앱)', '5,000,000원', '운영 자동화'),
        ('자부담', '5,000,000원', '대표자 자기 자금'),
    ]

    y = 1.4
    for i, (item, amount, note) in enumerate(budget_data_ppt):
        bg = PPTRGBColor(0xf8, 0xfa, 0xfc) if i % 2 == 0 else WHITE
        add_rect(slide7, 0.5, y, 12.3, 0.85, fill_color=bg, line_color=BLUE_LIGHT)
        add_text_box(slide7, item, 0.65, y + 0.15, 6, 0.6, font_size=12, bold=(i == 4), color=DARK)
        add_text_box(slide7, amount, 6.8, y + 0.15, 3, 0.6, font_size=12, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
        add_text_box(slide7, note, 10.0, y + 0.15, 2.7, 0.6, font_size=11, color=GRAY)
        y += 0.9

    add_rect(slide7, 0.5, y + 0.15, 12.3, 0.9, fill_color=BLUE)
    add_text_box(slide7, '합  계', 0.65, y + 0.25, 5, 0.65, font_size=14, bold=True, color=WHITE)
    add_text_box(slide7, '약 75,000,000원', 6.8, y + 0.25, 5.8, 0.65, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    # ── 슬라이드 8: 기대 효과 ──
    slide8 = add_slide()
    slide_header(slide8, '06. 기대 효과 및 성장 전략', '무인화 실현 → 연간 +3,000만 원 이상 순수익 증대')

    effects = [
        ('💸', '관리 인건비', '월 250만 원\n→ 0원', '+3,000만/년'),
        ('⚡', '전기료 절감', '30% 비용 절감', '+500만/년'),
        ('📋', '정산 자동화', '월 20시간 수작업\n→ 완전 자동', '기회비용 절감'),
        ('📈', '솔루션 확장', '외부 판매 목표\n2027년 SaaS 론칭', 'B2B 신규 수익'),
    ]
    for i, (icon, title, now, after) in enumerate(effects):
        x = 0.5 + i * 3.2
        add_rect(slide8, x, 1.4, 3.0, 5.5, fill_color=WHITE, line_color=BLUE_LIGHT)
        add_text_box(slide8, icon, x + 1.0, 1.55, 1.0, 0.7, font_size=28, align=PP_ALIGN.CENTER)
        add_text_box(slide8, title, x + 0.1, 2.3, 2.8, 0.55, font_size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text_box(slide8, now, x + 0.1, 2.95, 2.8, 1.3, font_size=11, color=GRAY, align=PP_ALIGN.CENTER)
        add_rect(slide8, x + 0.3, 4.4, 2.4, 0.7, fill_color=BLUE)
        add_text_box(slide8, after, x + 0.3, 4.45, 2.4, 0.6, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── 슬라이드 9: 마무리 ──
    slide9 = add_slide()
    add_rect(slide9, 0, 0, 13.33, 7.5, fill_color=DARK)
    add_rect(slide9, 0, 5.8, 13.33, 1.7, fill_color=BLUE)
    add_text_box(slide9, '"단순 임대업을 넘어\nAI 기반 무인 공유오피스의\n표준을 만든다"', 1, 1.0, 11, 3.5,
                 font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide9, '대우오피스  |  AI 무인 공유오피스 고도화 사업  |  2026', 1, 6.0, 11, 0.8,
                 font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

    output_path = os.path.join(OUTPUT_DIR, '대우오피스_사업계획서.pptx')
    prs.save(output_path)
    print(f'✅ PPT 저장 완료: {output_path}')
    return output_path


if __name__ == '__main__':
    print('📄 사업계획서 생성 시작...')
    word_path = create_word()
    ppt_path = create_ppt()
    print(f'\n✨ 생성 완료!')
    print(f'   Word: {word_path}')
    print(f'   PPT:  {ppt_path}')
